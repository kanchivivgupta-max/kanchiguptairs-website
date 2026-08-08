import os
import re
import sys
import subprocess
import json

def generate_notes_html(title, category, snippet, md_content, slug):
    html_lines = []
    paragraphs = md_content.strip().split("\n\n")
    
    for p in paragraphs:
        p = p.strip()
        if not p:
            continue
        if p.startswith("## "):
            html_lines.append(f"    <h2>{p[3:]}</h2>")
        elif p.startswith("### "):
            html_lines.append(f"    <h3>{p[4:]}</h3>")
        elif p.startswith("- ") or p.startswith("* "):
            list_items = []
            for line in p.split("\n"):
                line = line.strip()
                if line.startswith("- ") or line.startswith("* "):
                    line_content = line[2:]
                    line_content = re.sub(r"\*\*(.*?)\*\*", r"<strong>\1</strong>", line_content)
                    list_items.append(f"        <li>{line_content}</li>")
            html_lines.append("    <ul>\n" + "\n".join(list_items) + "\n    </ul>")
        else:
            p_content = re.sub(r"\*\*(.*?)\*\*", r"<strong>\1</strong>", p)
            p_content = p_content.replace("\n", "<br>")
            html_lines.append(f"    <p>{p_content}</p>")

    html_body = "\n".join(html_lines)

    template = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{title}: Technical Notes</title>
    <style>
    @page {{
        size: A4;
        margin: 2cm;
        @bottom-right {{
            content: "Page " counter(page) " of " counter(pages);
            font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
            font-size: 9pt;
            color: #666;
        }}
    }}
    body {{
        font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
        color: #24292e;
        line-height: 1.6;
        max-width: 900px;
        margin: 0 auto;
        padding: 40px 20px;
        background-color: #ffffff;
    }}
    h1, h2, h3, h4 {{
        color: #111;
        font-weight: 600;
        margin-top: 1.5em;
        margin-bottom: 0.5em;
        border-bottom: 1px solid #eaecef;
        padding-bottom: 0.3em;
    }}
    h1 {{ font-size: 26pt; color: #7C5CFF; border-bottom: 2px solid #7C5CFF; }}
    h2 {{ font-size: 20pt; color: #333; }}
    h3 {{ font-size: 15pt; border-bottom: none; color: #555; }}
    a {{ color: #7C5CFF; text-decoration: none; }}
    pre {{
        background-color: #f6f8fa;
        border-radius: 6px;
        padding: 16px;
        overflow: auto;
        font-family: 'Courier New', Courier, monospace;
        font-size: 9.5pt;
        border: 1px solid #e1e4e8;
    }}
    code {{
        font-family: 'Courier New', Courier, monospace;
        background-color: #f6f8fa;
        padding: 0.2em 0.4em;
        border-radius: 3px;
        font-size: 90%;
    }}
    pre code {{
        background-color: transparent;
        padding: 0;
    }}
    table {{
        border-collapse: collapse;
        width: 100%;
        margin: 20px 0;
    }}
    th, td {{
        border: 1px solid #dfe2e5;
        padding: 10px 14px;
        text-align: left;
    }}
    th {{
        background-color: #f6f8fa;
        font-weight: 600;
    }}
    tr:nth-child(2n) {{
        background-color: #f8f9fa;
    }}
    blockquote {{
        margin: 0;
        padding: 0 1em;
        color: #6a737d;
        border-left: 0.25em solid #dfe2e5;
    }}
    .callout {{
        margin: 20px 0;
        padding: 15px;
        border-radius: 6px;
        border-left: 5px solid #ccc;
    }}
    .callout-title {{
        font-weight: bold;
        margin-bottom: 8px;
        font-size: 11pt;
    }}
    .callout-important {{ background-color: #f4f0ff; border-left-color: #7C5CFF; color: #7C5CFF; }}
    .callout-important .callout-body {{ color: #24292e; }}
    .callout-warning {{ background-color: #fffbdd; border-left-color: #d73a49; color: #d73a49; }}
    .callout-warning .callout-body {{ color: #24292e; }}
    .callout-tip {{ background-color: #e6ffed; border-left-color: #28a745; color: #28a745; }}
    .callout-tip .callout-body {{ color: #24292e; }}
    .callout-note {{ background-color: #f6f8fa; border-left-color: #6a737d; color: #6a737d; }}
    .meta-info {{
        font-size: 10pt;
        color: #6a737d;
        margin-top: -10px;
        margin-bottom: 30px;
    }}
    </style>
</head>
<body>

    <h1>{title}</h1>
    <div class="meta-info">Compiled by Kanchi Gupta · Technical Study Series</div>

    <div class="callout callout-important">
        <div class="callout-title">Core Theme</div>
        <div class="callout-body">
            {snippet}
        </div>
    </div>

{html_body}

    <div style="margin-top: 50px; border-top: 1px solid #eaecef; padding-top: 20px; text-align: center;">
        <a href="../index.html">← Back to Homepage</a>
    </div>

</body>
</html>
"""
    return template

def generate_presentation_pptx(title, category, snippet, md_content, output_path):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("[PPTX] python-pptx not available. Skipping automatic PowerPoint generation.")
        return False

    print(f"[PPTX] Generating automated PowerPoint slide deck...")
    try:
        prs = Presentation()
        # Set to 16:9 landscape aspect ratio
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Color palette variables (Linear.app Style)
        color_bg = RGBColor(8, 9, 10)         # Marketing Black #08090a
        color_card = RGBColor(20, 21, 24)      # Surface Dark
        color_text_white = RGBColor(247, 248, 248) # Near white
        color_text_gray = RGBColor(138, 143, 152)  # Muted silver gray
        color_accent = RGBColor(113, 112, 255)     # Brand Violet #7170ff
        
        blank_slide_layout = prs.slide_layouts[6]
        
        # -------------------- SLIDE 1: COVER --------------------
        slide1 = prs.slides.add_slide(blank_slide_layout)
        background1 = slide1.background
        background1.fill.solid()
        background1.fill.fore_color.rgb = color_bg
        
        # Cover Text Box
        tx_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(4))
        tf = tx_box.text_frame
        tf.word_wrap = True
        
        # Main Title
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Arial"
        p_title.font.size = Pt(40)
        p_title.font.bold = True
        p_title.font.color.rgb = color_accent
        p_title.alignment = PP_ALIGN.CENTER
        p_title.space_after = Pt(20)
        
        # Subtitle
        p_sub = tf.add_paragraph()
        p_sub.text = "Professional Skilling & Adaptation in Public Service"
        p_sub.font.name = "Arial"
        p_sub.font.size = Pt(16)
        p_sub.font.color.rgb = color_text_gray
        p_sub.alignment = PP_ALIGN.CENTER
        p_sub.space_after = Pt(40)
        
        # Author details
        p_auth = tf.add_paragraph()
        p_auth.text = "Kanchi Gupta, Joint Commissioner (IRS)"
        p_auth.font.name = "Arial"
        p_auth.font.size = Pt(13)
        p_auth.font.bold = True
        p_auth.font.color.rgb = color_text_white
        p_auth.alignment = PP_ALIGN.CENTER
        
        p_meta = tf.add_paragraph()
        p_meta.text = "B.Tech (VNIT Nagpur) · M.Tech (IIT Bombay) · LL.B. (GLC Mumbai)"
        p_meta.font.name = "Courier New"
        p_meta.font.size = Pt(10)
        p_meta.font.color.rgb = color_text_gray
        p_meta.alignment = PP_ALIGN.CENTER

        # -------------------- SLIDES 2-4: CONTENT --------------------
        # Extract headers and content from Markdown
        headers = re.findall(r"^##\s*(.*?)$", md_content, re.MULTILINE)
        paragraphs = md_content.strip().split("\n\n")
        
        slide_count = 2
        for h in headers[:3]: # Limit to next 3 sections for 4-slide deck
            slide = prs.slides.add_slide(blank_slide_layout)
            background = slide.background
            background.fill.solid()
            background.fill.fore_color.rgb = color_bg
            
            # Slide Header
            header_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.75))
            htf = header_box.text_frame
            hp = htf.paragraphs[0]
            hp.text = f"The Lifelong Learning Loop   ·   Slide {slide_count} of 4"
            hp.font.name = "Arial"
            hp.font.size = Pt(9)
            hp.font.color.rgb = color_text_gray
            
            # Content Box
            content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(5))
            ctf = content_box.text_frame
            ctf.word_wrap = True
            
            # Title of Slide
            cp_title = ctf.paragraphs[0]
            cp_title.text = h
            cp_title.font.name = "Arial"
            cp_title.font.size = Pt(28)
            cp_title.font.bold = True
            cp_title.font.color.rgb = color_text_white
            cp_title.space_after = Pt(20)
            
            # Find paragraphs that are under this specific header
            # For simplicity, we grab the text that contains bullet points
            bullets_found = []
            for p in paragraphs:
                if p.strip().startswith("- ") or p.strip().startswith("* "):
                    # Extract bullet items
                    for line in p.split("\n"):
                        line = line.strip()
                        if line.startswith("- ") or line.startswith("* "):
                            # Remove bold markdown tags
                            clean_line = line[2:].replace("**", "")
                            bullets_found.append(clean_line)
                            
            # Add up to 3 bullets to Slide
            for bullet in bullets_found[:3]:
                bp = ctf.add_paragraph()
                bp.text = f"✦   {bullet}"
                bp.font.name = "Arial"
                bp.font.size = Pt(13)
                bp.font.color.rgb = color_text_gray
                bp.space_after = Pt(15)
                
            slide_count += 1
            
        prs.save(output_path)
        print(f"[PPTX] SUCCESS! PowerPoint presentation file compiled and saved: {output_path}")
        return True
    except Exception as e:
        print(f"[PPTX] Error compiling presentation slide deck: {e}")
        return False

def cloud_publish():
    path = os.getcwd()
    drafts_pool_dir = os.path.join(path, 'drafts_pool')
    log_file_path = os.path.join(path, 'published_log.txt')
    index_path = os.path.join(path, 'index.html')

    if not os.path.exists(drafts_pool_dir):
        print(f"Error: drafts_pool directory not found.")
        sys.exit(1)

    published_files = []
    if os.path.exists(log_file_path):
        with open(log_file_path, 'r', encoding='utf-8') as f:
            published_files = [line.strip() for line in f.read().split("\n") if line.strip()]

    all_drafts = sorted([f for f in os.listdir(drafts_pool_dir) if f.endswith('.md')])
    
    selected_draft_file = None
    for draft in all_drafts:
        if draft not in published_files:
            selected_draft_file = draft
            break

    if not selected_draft_file:
        print("[Cloud Content Factory] WARNING: All drafts published!")
        sys.exit(0)

    draft_full_path = os.path.join(drafts_pool_dir, selected_draft_file)
    print(f"[Cloud Content Factory] Active Draft Selected: {selected_draft_file}")

    with open(draft_full_path, 'r', encoding='utf-8') as f:
        content = f.read()

    frontmatter_match = re.match(r"^---\s*\n(.*?)\n---\s*\n(.*)$", content, re.DOTALL)
    if not frontmatter_match:
        print(f"Error: Draft {selected_draft_file} is missing YAML frontmatter.")
        sys.exit(1)

    frontmatter_raw = frontmatter_match.group(1)
    md_content = frontmatter_match.group(2)

    title = re.search(r"^title:\s*(.*?)$", frontmatter_raw, re.MULTILINE).group(1).strip()
    category = re.search(r"^category:\s*(.*?)$", frontmatter_raw, re.MULTILINE).group(1).strip()
    snippet = re.search(r"^snippet:\s*(.*?)$", frontmatter_raw, re.MULTILINE).group(1).strip()

    slug = title.replace(":", "").replace(" ", "_").replace("&", "and").strip()
    slug_url = f"notes/Master_{slug}.html"
    note_file_path = os.path.join(path, slug_url)

    print(f"[Cloud Content Factory] Generating HTML for: {title}")
    
    html_content = generate_notes_html(title, category, snippet, md_content, slug)

    with open(note_file_path, 'w', encoding='utf-8') as f:
        f.write(html_content)
    print("Success: Generated standalone HTML notes page.")

    # Generate PPTX PowerPoint Slide Deck automatically!
    pptx_path = os.path.join(path, f"notes/slides/{slug}_Presentation.pptx")
    os.makedirs(os.path.dirname(pptx_path), exist_ok=True)
    generate_presentation_pptx(title, category, snippet, md_content, pptx_path)

    # Update index.html
    with open(index_path, 'r', encoding='utf-8') as f:
        index_content = f.read()

    if slug_url in index_content:
        print("Note already linked on index.html. Skipping link insertion.")
    else:
        new_card = f"""        <!-- Note Card: {title} -->
        <div class="resource-card">
          <span class="resource-category">{category}</span>
          <h3>{title}</h3>
          <p>{snippet}</p>
          <div style="display: flex; gap: 1rem; margin-top: auto; padding-top: 1.25rem;">
            <a href="{slug_url}" target="_blank" class="resource-link">HTML Version →</a>
          </div>
        </div>

"""
        grid_pos = index_content.find('<div class="resource-grid">')
        if grid_pos == -1:
            print("Error: Could not find <div class=\"resource-grid\"> in index.html.")
            sys.exit(1)
            
        insert_index = index_content.find('\n', grid_pos) + 1
        updated_index_content = index_content[:insert_index] + new_card + index_content[insert_index:]
        
        with open(index_path, 'w', encoding='utf-8') as f:
            f.write(updated_index_content)
        print("Success: Injected new card link inside index.html's My Notes grid!")

    # Append to published log
    with open(log_file_path, 'a', encoding='utf-8') as f:
        f.write(selected_draft_file + "\n")
    print(f"Success: Marked {selected_draft_file} as published.")

    # Write sharing logs
    logs_dir = os.path.join(path, 'logs')
    os.makedirs(logs_dir, exist_ok=True)
    logs_file_path = os.path.join(logs_dir, 'scheduled_posts.txt')
    
    with open(logs_file_path, 'w', encoding='utf-8') as f:
        f.write(f"=============================================================================\n")
        f.write(f"SCHEDULED CONTENT FACTORY POST (CLOUD EXECUTED)\n")
        f.write(f"=============================================================================\n")
        f.write(f"Topic: {title}\n")
        f.write(f"URL: https://kanchiguptairs.com/{slug_url}\n")
        f.write(f"PPTX Slide Deck: https://kanchiguptairs.com/notes/slides/{slug}_Presentation.pptx\n")
        f.write(f"Category: {category}\n\n")
        f.write(f"--- DRAFT TEXT ---\n{md_content}\n")
    print("Success: Updated local log scheduled_posts.txt for sharing.")

if __name__ == '__main__':
    cloud_publish()
